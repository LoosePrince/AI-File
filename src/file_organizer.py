import os
import json
from PIL import Image
import openai
import base64
from io import BytesIO
import logging
from functools import reduce
import io
from PIL import Image
import time
import hashlib
import docx
import openpyxl
from pptx import Presentation
import requests.exceptions
from http.client import RemoteDisconnected
from urllib3.exceptions import ProtocolError
import configparser
import zipfile
import rarfile
import py7zr
import tempfile
import cv2
import subprocess
import datetime
from concurrent.futures import ThreadPoolExecutor
from threading import Lock
import re
import ctypes
import sys
import shutil
import stat

# 导入 Logger 类
from logger import Logger

# 导入子文件夹处理模式常量
from config import SUBFOLDER_MODE_WHOLE, SUBFOLDER_MODE_EXTRACT_ALL, SUBFOLDER_MODE_EXTRACT_PARTIAL

# 在文件顶部添加常量定义
class MessageType:
    ANALYZING = "分析中"
    CACHED_ANALYSIS = "使用缓存分析结果"
    CACHED_DECISION = "使用缓存的整理方案"
    ALL_CACHED_DECISION = "所有文件使用缓存的整理方案"
    GENERATING_DECISION = "生成整理方案中"
    GENERATING_DECISION_PROGRESS = "正在生成整理方案"
    DECISION_GENERATED = "整理方案生成完成"
    RETRY_BATCH = "重试当前批次"
    FILE_MOVED = "文件已移动到"
    FILE_ANALYZED = "文件已分析"
    FILE_NEED_ANALYSIS = "需要分析的文件"
    INFO = "info"
    WARNING = "warning"
    ERROR = "error"
    ANALYZING_FILE = "analyzing_file"

class FileOrganizer:
    def __init__(self, api_key, output_dir=None):
        self.api_key = api_key
        self.output_dir = output_dir
        self.cancel_flag = False
        self.progress_callback = None
        self.progress_lock = Lock()
        self.analysis_cache = {}
        self.decision_cache = {}
        self.logger = Logger()
        self.prompt = None  # 添加提示词属性
        self.load_config()

    def load_config(self):
        """加载或更新配置"""
        try:
            from config import API_KEY, API_URL, API_TYPE, FILE_OPERATION, IMAGE_MODEL, FILE_MODEL, DECISION_MODEL, ENABLE_VIDEO_ANALYSIS, VIDEO_MODEL, LANGUAGE, SUBFOLDER_MODE
            
            self.logger.log_info("配置文件加载成功")
            
            # 获取文件操作模式
            self.file_operation = FILE_OPERATION
            # 获取模型设置
            self.image_model = IMAGE_MODEL
            self.file_model = FILE_MODEL
            self.decision_model = DECISION_MODEL
            # 获取视频分析设置
            self.enable_video_analysis = ENABLE_VIDEO_ANALYSIS
            self.video_model = VIDEO_MODEL
            # 获取语言设置
            self.language = LANGUAGE
            # 获取子文件夹处理模式
            self.subfolder_mode = SUBFOLDER_MODE
            # 获取API类型
            self.api_type = API_TYPE
            
            # 更新openai设置
            openai.api_key = API_KEY
            openai.api_base = API_URL
            
        except Exception as e:
            self.logger.log_error(f"加载配置文件时出错: {str(e)}")
            # 使用默认值
            self.file_operation = 'copy'
            self.image_model = 'Pro/Qwen/Qwen2-VL-7B-Instruct'
            self.file_model = 'deepseek-ai/DeepSeek-R1-Distill-Qwen-7B'
            self.decision_model = 'deepseek-ai/DeepSeek-R1-Distill-Qwen-32B'
            self.enable_video_analysis = False
            self.video_model = 'Pro/Qwen/Qwen2-VL-7B-Instruct'
            self.language = '中文'
            self.subfolder_mode = 'whole'  # 默认整体处理
            self.api_type = 'OpenAI API'  # 默认使用 OpenAI API

    def set_progress_callback(self, callback):
        self.progress_callback = callback

    def set_prompt(self, prompt):
        """设置提示词"""
        self.prompt = prompt

    def _load_history(self):
        if os.path.exists(self.history_file):
            try:
                with open(self.history_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except PermissionError:
                if self._handle_permission_error(self.history_file, "读取"):
                    try:
                        with open(self.history_file, 'r', encoding='utf-8') as f:
                            return json.load(f)
                    except Exception as e:
                        self.logger.log_error(f"读取历史记录失败：{str(e)}")
                        return {}
                return {}
            except Exception as e:
                self.logger.log_error(f"读取历史记录失败：{str(e)}")
                return {}
        return {}

    def _load_analysis_cache(self):
        if os.path.exists(self.analysis_cache_file):
            try:
                with open(self.analysis_cache_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except PermissionError:
                if self._handle_permission_error(self.analysis_cache_file, "读取"):
                    try:
                        with open(self.analysis_cache_file, 'r', encoding='utf-8') as f:
                            return json.load(f)
                    except Exception as e:
                        self.logger.log_error(f"读取分析缓存失败：{str(e)}")
                        return {}
                return {}
            except Exception as e:
                self.logger.log_error(f"读取分析缓存失败：{str(e)}")
                return {}
        return {}

    def _load_decision_cache(self):
        if os.path.exists(self.decision_cache_file):
            try:
                with open(self.decision_cache_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except PermissionError:
                if self._handle_permission_error(self.decision_cache_file, "读取"):
                    try:
                        with open(self.decision_cache_file, 'r', encoding='utf-8') as f:
                            return json.load(f)
                    except Exception as e:
                        self.logger.log_error(f"读取决策缓存失败：{str(e)}")
                        return {}
                return {}
            except Exception as e:
                self.logger.log_error(f"读取决策缓存失败：{str(e)}")
                return {}
        return {}

    def _handle_permission_error(self, file_path, operation="访问"):
        """处理权限不足的情况，尝试获取权限"""
        try:
            if sys.platform == 'win32':
                # Windows系统下尝试获取管理员权限
                if not ctypes.windll.shell32.IsUserAnAdmin():
                    self.logger.log_warning(f"尝试获取管理员权限以{operation}文件：{file_path}")
                    # 重新以管理员权限运行程序
                    ctypes.windll.shell32.ShellExecuteW(None, "runas", sys.executable, " ".join(sys.argv), None, 1)
                    return False
            else:
                # Linux/Unix系统下尝试修改文件权限
                current_permissions = os.stat(file_path).st_mode
                # 添加读写权限
                os.chmod(file_path, current_permissions | stat.S_IRUSR | stat.S_IWUSR)
                self.logger.log_info(f"已修改文件权限：{file_path}")
            
            return True
        except Exception as e:
            self.logger.log_error(f"获取文件权限失败：{str(e)}")
            return False

    def _safe_file_operation(self, operation, source_path, target_path=None):
        """安全地执行文件操作，处理权限问题"""
        try:
            if operation == "move":
                if target_path:
                    # 确保目标目录存在
                    os.makedirs(os.path.dirname(target_path), exist_ok=True)
                # 使用 shutil.move 时添加 dirs_exist_ok=True 参数
                shutil.move(source_path, target_path)
            elif operation == "copy":
                if target_path:
                    # 确保目标目录存在
                    os.makedirs(os.path.dirname(target_path), exist_ok=True)
                # 使用 shutil.copy2 时添加 dirs_exist_ok=True 参数
                shutil.copy2(source_path, target_path)
            elif operation == "delete":
                # 如果是文件夹，使用 shutil.rmtree
                if os.path.isdir(source_path):
                    shutil.rmtree(source_path)
                else:
                    os.remove(source_path)
            elif operation == "mkdir":
                # 使用 exist_ok=True 参数，避免文件夹已存在时的错误
                os.makedirs(source_path, exist_ok=True)
            return True
        except PermissionError:
            # 遇到权限错误时尝试获取权限
            if self._handle_permission_error(source_path, operation):
                try:
                    # 重试操作
                    if operation == "move":
                        shutil.move(source_path, target_path)
                    elif operation == "copy":
                        shutil.copy2(source_path, target_path)
                    elif operation == "delete":
                        if os.path.isdir(source_path):
                            shutil.rmtree(source_path)
                        else:
                            os.remove(source_path)
                    elif operation == "mkdir":
                        os.makedirs(source_path, exist_ok=True)
                    return True
                except Exception as e:
                    self.logger.log_error(f"重试文件操作失败：{str(e)}")
            return False
        except Exception as e:
            self.logger.log_error(f"文件操作失败：{str(e)}")
            return False

    def _save_history(self):
        try:
            with open(self.history_file, 'w', encoding='utf-8') as f:
                json.dump(self.history, f, indent=2, ensure_ascii=False)
        except PermissionError:
            if self._handle_permission_error(self.history_file, "写入"):
                try:
                    with open(self.history_file, 'w', encoding='utf-8') as f:
                        json.dump(self.history, f, indent=2, ensure_ascii=False)
                except Exception as e:
                    self.logger.log_error(f"保存历史记录失败：{str(e)}")
        except Exception as e:
            self.logger.log_error(f"保存历史记录失败：{str(e)}")

    def _save_analysis_cache(self):
        try:
            with open(self.analysis_cache_file, 'w', encoding='utf-8') as f:
                json.dump(self.analysis_cache, f, indent=2, ensure_ascii=False)
        except PermissionError:
            if self._handle_permission_error(self.analysis_cache_file, "写入"):
                try:
                    with open(self.analysis_cache_file, 'w', encoding='utf-8') as f:
                        json.dump(self.analysis_cache, f, indent=2, ensure_ascii=False)
                except Exception as e:
                    self.logger.log_error(f"保存分析缓存失败：{str(e)}")
        except Exception as e:
            self.logger.log_error(f"保存分析缓存失败：{str(e)}")

    def _save_decision_cache(self):
        try:
            with open(self.decision_cache_file, 'w', encoding='utf-8') as f:
                json.dump(self.decision_cache, f, indent=2, ensure_ascii=False)
        except PermissionError:
            if self._handle_permission_error(self.decision_cache_file, "写入"):
                try:
                    with open(self.decision_cache_file, 'w', encoding='utf-8') as f:
                        json.dump(self.decision_cache, f, indent=2, ensure_ascii=False)
                except Exception as e:
                    self.logger.log_error(f"保存决策缓存失败：{str(e)}")
        except Exception as e:
            self.logger.log_error(f"保存决策缓存失败：{str(e)}")

    def _calculate_md5(self, file_path):
        
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()

    def _get_directory_structure(self, rootdir):
        """获取目录结构"""
        dir_structure = {}
        rootdir = rootdir.rstrip(os.sep)
        start = rootdir.rfind(os.sep) + 1
        
        try:
            for path, dirs, files in os.walk(rootdir):
                # 跳过系统文件夹和隐藏文件夹
                dirs[:] = [d for d in dirs if not d.startswith('.') and not d.startswith('$')]
                files = [f for f in files if not f.startswith('.') and not f.startswith('$')]
                
                folders = path[start:].split(os.sep)
                subdir = dict.fromkeys(files)
                parent = reduce(dict.get, folders[:-1], dir_structure)
                parent[folders[-1]] = subdir
        except Exception as e:
            self.logger.log_error(f"获取目录结构失败：{str(e)}")
            return {}
            
        return dir_structure

    def _call_api_with_retry(self, api_func, **kwargs):
        """带重试机制的API调用"""
        max_retries = 3
        retry_count = 0
        last_error = None
        
        while retry_count < max_retries:
            try:
                if self.api_type == "Ollama API":
                    # 修改消息格式以适应 Ollama API
                    if "messages" in kwargs:
                        messages = kwargs["messages"]
                        # 将消息转换为 Ollama 格式
                        prompt = ""
                        for msg in messages:
                            if msg["role"] == "system":
                                prompt += f"System: {msg['content']}\n"
                            elif msg["role"] == "user":
                                prompt += f"User: {msg['content']}\n"
                            elif msg["role"] == "assistant":
                                prompt += f"Assistant: {msg['content']}\n"
                        
                        # 调用 Ollama API
                        import requests
                        response = requests.post(
                            f"{openai.api_base}/api/generate",
                            json={
                                "model": kwargs.get("model", "llama2"),
                                "prompt": prompt,
                                "stream": False
                            }
                        )
                        
                        if response.status_code == 200:
                            result = response.json()
                            # 构造类似 OpenAI 的响应格式
                            return type('Response', (), {
                                'choices': [
                                    type('Choice', (), {
                                        'message': type('Message', (), {
                                            'content': result.get('response', '')
                                        })
                                    })
                                ]
                            })
                        else:
                            raise Exception(f"Ollama API 调用失败: {response.text}")
                    else:
                        raise Exception("Ollama API 不支持当前调用方式")
                else:
                    # 使用 OpenAI API
                    return api_func(**kwargs)
                    
            except Exception as e:
                last_error = e
                retry_count += 1
                if retry_count < max_retries:
                    time.sleep(1)  # 等待1秒后重试
                continue
                
        raise last_error or Exception("API调用失败")

    def _compress_image(self, img, quality=50):
        """压缩图像"""
        output = BytesIO()
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        img.save(output, format='JPEG', quality=quality, optimize=True)
        output.seek(0)
        return Image.open(output)

    def _move_or_copy_file(self, src, dst):
        """根据设置移动或复制文件"""
        try:
            # 确保目标目录存在
            os.makedirs(os.path.dirname(dst), exist_ok=True)
            
            # 判断源路径是文件还是文件夹
            if os.path.isdir(src):
                if self.file_operation == 'copy':
                    # 复制整个文件夹及其内容，使用 dirs_exist_ok=True
                    shutil.copytree(src, dst)
                    if self.progress_callback:
                        self.progress_callback(f"复制文件夹：{src} -> {dst}")
                else:
                    # 移动整个文件夹，使用 dirs_exist_ok=True
                    shutil.move(src, dst)
                    if self.progress_callback:
                        self.progress_callback(f"移动文件夹：{src} -> {dst}")
            else:
                # 处理单个文件
                if self.file_operation == 'copy':
                    shutil.copy2(src, dst)  # copy2保留文件的元数据
                    if self.progress_callback:
                        self.progress_callback(f"复制文件：{src} -> {dst}")
                else:
                    # 移动文件时使用 shutil.move，添加 dirs_exist_ok=True
                    shutil.move(src, dst)
                    if self.progress_callback:
                        self.progress_callback(f"移动文件：{src} -> {dst}")
            return True
        except Exception as e:
            op_type = '复制' if self.file_operation == 'copy' else '移动'
            file_type = '文件夹' if os.path.isdir(src) else '文件'
            if self.progress_callback:
                self.progress_callback(f"{op_type}{file_type}失败：{src} -> {dst}，错误：{str(e)}")
            print(f"{op_type}{file_type}失败：{src} -> {dst}，错误：{str(e)}")
            return False

    def _analyze_image(self, image_path):
        try:
            # 检查缓存
            file_md5 = self._calculate_md5(image_path)
            if file_md5 in self.analysis_cache:
                if self.progress_callback:
                    self.progress_callback(MessageType.CACHED_ANALYSIS)
                return self.analysis_cache[file_md5]
            
            if self.progress_callback:
                self.progress_callback(MessageType.ANALYZING)
            
            # 检查取消标志
            if self.cancel_flag:
                if self.progress_callback:
                    self.progress_callback("操作已取消")
                return None
            
            # 打开并检查图像大小
            img = Image.open(image_path)
            file_size = os.path.getsize(image_path) / (1024 * 1024)  # 转换为MB
            compression_ratio = 1.0
            compressed_size = file_size
            
            # 根据文件大小决定压缩率
            if file_size > 1:
                # 计算需要的压缩率以达到1MB
                target_size = 1  # 目标大小1MB
                needed_ratio = target_size / file_size
                # 限制最大压缩率为85%
                compression_ratio = max(0.15, needed_ratio)  # 0.15表示压缩85%
                quality = int(compression_ratio * 100)  # 转换为quality参数
                
                if self.progress_callback:
                    compress_percent = int((1 - compression_ratio) * 100)
                    self.progress_callback(f"正在压缩图像（{compress_percent}%）：{image_path}")
                
                # 检查取消标志
                if self.cancel_flag:
                    if self.progress_callback:
                        self.progress_callback("操作已取消")
                    return None
                    
                img = self._compress_image(img, quality=quality)
            
            # 计算压缩后的大小
            buffered = BytesIO()
            img.save(buffered, format="JPEG" if img.format == "JPEG" else "PNG")
            compressed_size = len(buffered.getvalue()) / (1024 * 1024)  # 转换为MB
            img_str = base64.b64encode(buffered.getvalue()).decode('utf-8')
            
            # 检查取消标志
            if self.cancel_flag:
                if self.progress_callback:
                    self.progress_callback("操作已取消")
                return None
            
            # 使用重试机制调用API
            response = self._call_api_with_retry(
                openai.ChatCompletion.create,
                model=self.image_model,
                messages=[{
                    "role": "system",
                    "content": "使用不超过50字的一句话描述分析图像内容并描述其主要特征（如画风、软件、主题、人物）和概要可能包含的信息"
                },{
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{img_str}"
                            }
                        },
                        {
                            "type": "text",
                            "text": "请分析这张图片的内容"
                        }
                    ]
                }]
            )
            
            # 检查取消标志
            if self.cancel_flag:
                if self.progress_callback:
                    self.progress_callback("操作已取消")
                return None
            
            result = {
                "type": "image",
                "content": response.choices[0].message.content,
                "metadata": {
                    "format": img.format,
                    "size": img.size,
                    "mode": img.mode,
                    "original_size_mb": file_size,
                    "compressed": file_size > 1,
                    "compression_ratio": compression_ratio,
                    "compressed_size_mb": compressed_size
                }
            }
            
            # 更新缓存
            self.analysis_cache[file_md5] = result
            self._save_analysis_cache()
            return result
            
        except Exception as e:
            self.logger.log_error(f"分析图像 {image_path} 失败：{str(e)}")
            raise Exception(f"分析图像 {image_path} 失败：{str(e)}")

    def _extract_docx_content(self, file_path):
        """提取Word文档内容"""
        try:
            
            doc = docx.Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
            # 限制文本长度，避免过长
            if len(text) > 1000:
                text = text[:1000] + "..."
            return text
        except Exception as e:
            print(f"提取Word文档内容失败：{str(e)}")
            return ""
    
    def _extract_xlsx_content(self, file_path):
        """提取Excel文件内容"""
        try:
            
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text_parts = []
            
            # 最多处理前3个工作表
            sheet_count = 0
            for sheet_name in wb.sheetnames:
                if sheet_count >= 3:
                    break
                    
                sheet = wb[sheet_name]
                sheet_text = f"工作表：{sheet_name}\n"
                
                # 最多读取每个工作表的前50行、10列数据
                row_count = 0
                for row in sheet.iter_rows(max_row=50, max_col=10, values_only=True):
                    row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        sheet_text += row_text + "\n"
                    row_count += 1
                    if row_count >= 50:
                        break
                        
                text_parts.append(sheet_text)
                sheet_count += 1
                
            text = "\n\n".join(text_parts)
            # 限制文本长度
            if len(text) > 1000:
                text = text[:1000] + "..."
            return text
        except Exception as e:
            print(f"提取Excel文件内容失败：{str(e)}")
            return ""
    
    def _extract_pptx_content(self, file_path):
        """提取PowerPoint文件内容"""
        try:
            
            prs = Presentation(file_path)
            text_parts = []
            
            # 最多处理前10张幻灯片
            slide_count = 0
            for slide in prs.slides:
                if slide_count >= 10:
                    break
                    
                slide_text = f"幻灯片 {slide_count+1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                        
                if slide_text.strip() != f"幻灯片 {slide_count+1}:":
                    text_parts.append(slide_text)
                slide_count += 1
                
            text = "\n\n".join(text_parts)
            # 限制文本长度
            if len(text) > 1000:
                text = text[:1000] + "..."
            return text
        except Exception as e:
            print(f"提取PowerPoint文件内容失败：{str(e)}")
            return ""
    
    def _extract_archive_structure(self, file_path, max_depth=4, max_entries=50):
        """提取压缩包的目录结构"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            structure = {"files": [], "directories": []}
            entries_count = 0
            
            if ext in ['.zip']:
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    namelist = zip_ref.namelist()
                    for name in namelist:
                        if entries_count >= max_entries:
                            break
                        # 计算当前路径的深度
                        depth = name.count('/') + (1 if name.endswith('/') else 0)
                        if depth > max_depth:
                            continue
                        if name.endswith('/'):
                            structure["directories"].append(name)
                        else:
                            structure["files"].append(name)
                        entries_count += 1
                        
            elif ext in ['.rar']:
                with rarfile.RarFile(file_path, 'r') as rar_ref:
                    infolist = rar_ref.infolist()
                    for info in infolist:
                        if entries_count >= max_entries:
                            break
                        # 计算当前路径的深度
                        depth = info.filename.count('/') + (1 if info.is_dir() else 0)
                        if depth > max_depth:
                            continue
                        if info.is_dir():
                            structure["directories"].append(info.filename)
                        else:
                            structure["files"].append(info.filename)
                        entries_count += 1
                        
            elif ext in ['.7z']:
                with py7zr.SevenZipFile(file_path, 'r') as sz_ref:
                    for filename, _ in sz_ref.list():
                        if entries_count >= max_entries:
                            break
                        # 计算当前路径的深度
                        depth = filename.count('/') + (1 if filename.endswith('/') else 0)
                        if depth > max_depth:
                            continue
                        if filename.endswith('/'):
                            structure["directories"].append(filename)
                        else:
                            structure["files"].append(filename)
                        entries_count += 1
                        
            return structure
        except Exception as e:
            print(f"提取压缩包结构失败：{str(e)}")
            return None
            
    def _analyze_archive(self, file_path):
        """分析压缩包内容"""
        try:
            # 提取压缩包结构
            archive_structure = self._extract_archive_structure(file_path)
            if not archive_structure:
                return None
                
            # 构建提示信息
            filename = os.path.basename(file_path)
            ext = os.path.splitext(filename)[1].lower()
            archive_type = "压缩包"
            
            # 格式化目录结构
            structure_text = "目录结构：\n"
            if archive_structure["directories"]:
                structure_text += "\n文件夹：\n"
                for dir_path in sorted(archive_structure["directories"])[:20]:
                    structure_text += f"- {dir_path}\n"
            if archive_structure["files"]:
                structure_text += "\n文件：\n"
                for file_path in sorted(archive_structure["files"])[:30]:
                    structure_text += f"- {file_path}\n"
            
            prompt = f"""文件名：{filename}
文件类型：{archive_type}
{structure_text}"""
            
            # 使用重试机制调用API
            response = self._call_api_with_retry(
                openai.ChatCompletion.create,
                model=self.file_model,
                messages=[{
                    "role": "system",
                    "content": "根据压缩包的文件名和内部目录结构，使用不超过50字的一句话描述压缩包的主要内容和用途"
                },{
                    "role": "user",
                    "content": prompt
                }]
            )
            
            result = {
                "type": "archive",
                "content": response.choices[0].message.content,
                "metadata": {
                    "filename": filename,
                    "extension": ext,
                    "file_type": archive_type,
                    "structure": archive_structure
                }
            }
            
            return result
            
        except Exception as e:
            raise Exception(f"分析压缩包 {file_path} 失败：{str(e)}")
            
    def _extract_text_content(self, file_path, max_chars=2000):
        """提取文本文件内容
        Args:
            file_path: 文件路径
            max_chars: 最大读取字符数，默认2000字符
        Returns:
            str: 提取的文本内容
        """
        try:
            # 尝试不同的编码方式读取文件
            encodings = ['utf-8', 'gbk', 'gb2312', 'ascii']
            content = ""
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read(max_chars)
                        if len(content) == max_chars:
                            content += "..."
                        break
                except UnicodeDecodeError:
                    continue
                except Exception as e:
                    self.logger.log_error(f"使用 {encoding} 编码读取文件 {file_path} 失败: {str(e)}")
                    continue
                
            if not content:
                self.logger.log_warning(f"无法读取文件内容: {file_path}")
                return ""
            
            return content.strip()
        
        except Exception as e:
            self.logger.log_error(f"提取文本内容失败: {str(e)}")
            return ""

    def _analyze_file(self, file_path):
        # 检查缓存
        file_md5 = self._calculate_md5(file_path)
        if file_md5 in self.analysis_cache:
            if self.progress_callback:
                self.progress_callback(MessageType.CACHED_ANALYSIS)
            return self.analysis_cache[file_md5]
            
        if self.progress_callback:
            self.progress_callback(MessageType.ANALYZING)
            
        # 检查取消标志
        if self.cancel_flag:
            if self.progress_callback:
                self.progress_callback("操作已取消")
            return None
            
        try:
            filename = os.path.basename(file_path)
            ext = os.path.splitext(filename)[1].lower()
            
            # 根据文件类型选择分析方法
            if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.tif']:
                # 对图像文件使用专门的分析方法
                result = self._analyze_image(file_path)
                return result
            elif ext in ['.zip', '.rar', '.7z']:
                result = self._analyze_archive(file_path)
            else:
                # 提取文件内容
                file_content = ""
                file_type = "普通文件"
                
                # 根据文件类型提取内容
                if ext == ".docx":
                    file_content = self._extract_docx_content(file_path)
                    file_type = "Word文档"
                elif ext == ".xlsx" or ext == ".xls":
                    file_content = self._extract_xlsx_content(file_path)
                    file_type = "Excel表格"
                elif ext == ".pptx" or ext == ".ppt":
                    file_content = self._extract_pptx_content(file_path)
                    file_type = "PowerPoint演示文稿"
                # 添加文本文件处理
                elif ext in ['.txt', '.md', '.log', '.ini', '.json', '.xml', '.csv', '.yml', '.yaml', '.conf']:
                    file_content = self._extract_text_content(file_path)
                    if ext == '.md':
                        file_type = "Markdown文档"
                    elif ext == '.log':
                        file_type = "日志文件"
                    elif ext in ['.json', '.xml', '.yml', '.yaml']:
                        file_type = "配置/数据文件"
                    elif ext == '.csv':
                        file_type = "CSV数据文件"
                    else:
                        file_type = "文本文件"
                
                # 检查取消标志
                if self.cancel_flag:
                    if self.progress_callback:
                        self.progress_callback("操作已取消")
                    return None
                
                # 构建提示信息
                prompt = f"文件名：{filename}\n文件类型：{file_type}"
                if file_content:
                    prompt += f"\n文件内容摘要：\n{file_content}"
                
                # 使用重试机制调用API
                response = self._call_api_with_retry(
                    openai.ChatCompletion.create,
                    model=self.file_model,
                    messages=[{
                        "role": "system",
                        "content": "根据文件名、文件类型和文件内容使用不超过50字的一句话描述分析文件内容"
                    },{
                        "role": "user",
                        "content": prompt
                    }]
                )
                
                # 检查取消标志
                if self.cancel_flag:
                    if self.progress_callback:
                        self.progress_callback("操作已取消")
                    return None
                
                result = {
                    "type": "file",
                    "content": response.choices[0].message.content,
                    "metadata": {
                        "filename": filename,
                        "extension": ext,
                        "file_type": file_type
                    }
                }
            
            # 更新缓存
            if result:
                self.analysis_cache[file_md5] = result
                self._save_analysis_cache()
            return result
            
        except Exception as e:
            raise Exception(f"分析文件 {file_path} 失败：{str(e)}")

    def _get_final_decision(self, analysis_results):
        max_retries = 3
        retry_count = 0
        all_files = list(analysis_results.items())
        total_files = len(all_files)
        batch_size = 10
        decisions = []
        
        # 用于存储已创建的文件夹及其类型映射
        folder_categories = {}
        
        if self.progress_callback:
            self.progress_callback(MessageType.GENERATING_DECISION)
            
        # 检查缓存中是否已有部分或全部结果
        cached_decisions = []
        uncached_files = []
        cached_count = 0
        
        # 首先从缓存的决策中收集已有的文件夹类型
        for file_path, analysis in all_files:
            if self.cancel_flag:  # 检查取消标志
                if self.progress_callback:
                    self.progress_callback("操作已取消")
                return None
                
            # 检查是否是文件夹
            is_folder = os.path.isdir(file_path)
            if is_folder:
                # 对于文件夹，直接使用文件夹分析结果
                folder_analysis = analysis
                if folder_analysis and isinstance(folder_analysis, dict):
                    # 使用文件夹分析结果生成决策
                    category = folder_analysis.get('category', '其他')
                    new_path = os.path.join(category, os.path.basename(file_path))
                    decisions.append({
                        'original_path': file_path,
                        'new_path': new_path,
                        'is_folder': True,
                        'analysis': folder_analysis,
                        'category': category,
                        'purpose': folder_analysis.get('purpose', '未知'),
                        'description': folder_analysis.get('description', f"包含{folder_analysis.get('file_count', 0)}个文件的文件夹")
                    })
                    
                    # 显示文件夹的目标路径
                    if self.progress_callback:
                        target_path = os.path.join(self.output_dir, new_path.lstrip(os.sep)) if self.output_dir else new_path
                        self.progress_callback(f"文件夹 {file_path} 将移动到：{target_path}")
                    continue
            
            file_md5 = self._calculate_md5(file_path)
            if file_md5 in self.decision_cache:
                cached_info = self.decision_cache[file_md5]
                
                # 检查文件是否存在且在正确的目录中
                if not os.path.exists(file_path):
                    if self.progress_callback:
                        self.progress_callback(f"跳过缓存：文件不存在 - {file_path}")
                    uncached_files.append((file_path, analysis))
                    continue
                
                # 验证文件是否在当前选择的目录中
                normalized_file_path = os.path.normpath(file_path)
                normalized_dir_path = os.path.normpath(os.path.dirname(file_path))
                
                # 如果指定了输出目录，确保新路径在输出目录中
                if self.output_dir:
                    new_path = os.path.join(self.output_dir, cached_info["new_path"].lstrip(os.sep))
                    if not self._validate_output_path(new_path):
                        if self.progress_callback:
                            self.progress_callback(f"跳过缓存：目标路径不在输出目录中 - {new_path}")
                        uncached_files.append((file_path, analysis))
                        continue
                
                # 使用缓存的决策结果
                cached_decision = cached_info.copy()
                # 更新原始路径，因为同一个文件可能在不同位置
                cached_decision["original_path"] = file_path
                cached_decisions.append(cached_decision)
                
                # 记录文件夹类型
                folder_name = os.path.dirname(cached_decision["new_path"])
                if folder_name:
                    folder_categories[folder_name.lower()] = folder_name
                
                cached_count += 1
            else:
                uncached_files.append((file_path, analysis))
        
        # 如果所有文件都有缓存结果，直接返回
        if cached_count == total_files:
            if self.progress_callback:
                self.progress_callback(MessageType.ALL_CACHED_DECISION)
            return {"files": cached_decisions}
        
        # 处理未缓存的文件
        all_files = uncached_files
        total_uncached = len(all_files)
        
        if self.progress_callback:
            self.progress_callback(f"开始分析 {total_uncached} 个未缓存文件...")
        
        for i in range(0, total_uncached, batch_size):
            batch = all_files[i:i+batch_size]
            batch_retry_count = 0
            current_batch = i // batch_size + 1
            total_batches = (total_uncached + batch_size - 1) // batch_size
            
            # 计算当前批次开始前的进度百分比
            progress_base = int((cached_count / total_files) * 100)
            current_progress = progress_base + int(((current_batch - 1) / total_batches) * (100 - progress_base))
            
            # 显示当前批次的文件列表
            if self.progress_callback:
                batch_files = "\n".join([f"- {file_path}" for file_path, _ in batch])
                self.progress_callback(f"正在分析第 {current_batch}/{total_batches} 批文件：\n{batch_files}")
            
            while batch_retry_count < max_retries:
                try:
                    if self.progress_callback:
                        self.progress_callback(f"正在生成整理方案（第{current_batch}/{total_batches}批，共{total_uncached}个未缓存文件）", current_progress)

                    # 使用配置的决策模型
                    lang_prompt = "使用中文" if self.language == "中文" else "使用" + self.language
                    output_dir = f"新路径在{self.output_dir}文件夹内整理" if self.output_dir else "新路径使用相对路径"
                    dir_structure = self._get_directory_structure(self.output_dir) if self.output_dir else "无指定输出目录"
                    
                    # 添加已有文件夹类型信息到提示中
                    existing_folders = "已有的文件夹类型：" + ", ".join(folder_categories.values()) if folder_categories else "暂无已创建的文件夹"
                    
                    # 构建系统提示词
                    system_prompt = """
请严格按照以下JSON格式输出（注意：不要输出多余的内容，只输出json，不要使用代码块包裹）：
{
  "files": [
    {
      "original_path": "原文件路径", 
      "type": "文件类型",
      "description": "文件描述",
      "new_path": "新路径"
    }
  ]
}"""
                    
                    # 构建用户提示词
                    user_prompt = f"""根据分析结果和当前目录结构，为每批文件建议最佳的组织方式。
当前输出目录结构：{dir_structure}
{existing_folders}
请优先使用已有的文件夹类型，避免创建相似功能的重复文件夹（如"压缩包"和"压缩文件"）。
{lang_prompt}命名文件夹，{output_dir}，输出的路径需要包含文件名。"""

                    # 如果有自定义提示词，添加到用户提示词中
                    if self.prompt:
                        user_prompt += f"\n\n用户自定义要求：\n{self.prompt}"

                    batch_data = {k: v for k, v in batch}
                    response = self._call_api_with_retry(
                        openai.ChatCompletion.create,
                        model=self.decision_model,
                        messages=[
                            {"role": "user", "content": user_prompt},
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": f"当前批次分析结果（共{len(batch)}个文件）：\n{json.dumps(batch_data, indent=2)}"}
                        ]
                    )

                    content = response.choices[0].message.content.strip()
                    if content.startswith("```json") and content.endswith("```"):
                        content = content[7:-3].strip()
                    decision = json.loads(content)
                    
                    if not isinstance(decision, dict) or "files" not in decision:
                        raise ValueError("决策格式错误：缺少files字段")
                    
                    # 更新决策缓存和文件夹类型映射
                    for file_info in decision["files"]:
                        original_path = file_info["original_path"]
                        
                        # 验证文件是否存在
                        if not os.path.exists(original_path):
                            continue
                            
                        # 如果指定了输出目录，验证新路径
                        if self.output_dir:
                            new_path = os.path.join(self.output_dir, file_info["new_path"].lstrip(os.sep))
                            if not self._validate_output_path(new_path):
                                continue
                        
                        file_md5 = self._calculate_md5(original_path)
                        # 复制一份进入缓存，移除原始路径，因为路径可能变化
                        cache_info = file_info.copy()
                        self.decision_cache[file_md5] = cache_info
                        
                        # 更新文件夹类型映射
                        folder_name = os.path.dirname(file_info["new_path"])
                        if folder_name:
                            folder_categories[folder_name.lower()] = folder_name
                            
                        # 显示文件的目标路径
                        if self.progress_callback:
                            target_path = os.path.join(self.output_dir, file_info["new_path"].lstrip(os.sep)) if self.output_dir else file_info["new_path"]
                            self.progress_callback(f"文件 {original_path} 将移动到：{target_path}")
                    
                    decisions.extend(decision["files"])
                    
                    break
                except Exception as e:
                    self.logger.log_error(f"生成决策失败：{str(e)}")
                    if retry_count < max_retries:
                        retry_count += 1
                        if self.progress_callback:
                            self.progress_callback(MessageType.RETRY_BATCH)
                        continue
                    else:
                        if self.progress_callback:
                            self.progress_callback(f"生成决策失败：{str(e)}")
                        return None
        
        # 合并所有决策
        final_decision = {
            'files': cached_decisions + decisions
        }
        
        if self.progress_callback:
            self.progress_callback(MessageType.DECISION_GENERATED)
            
        return final_decision

    def _validate_output_path(self, path):
        """验证输出路径是否在允许的输出目录中"""
        if not self.output_dir:
            # 如果没有指定输出目录，则使用相对路径，此时允许在源文件所在目录创建子目录
            return True
            
        # 规范化路径
        normalized_path = os.path.normpath(path)
        normalized_output_dir = os.path.normpath(self.output_dir)
        
        # 检查路径是否在输出目录中
        try:
            # 使用os.path.commonpath检查路径是否在输出目录下
            common_path = os.path.commonpath([normalized_path, normalized_output_dir])
            return common_path == normalized_output_dir
        except ValueError:
            # 如果路径不在同一个驱动器或有其他问题，返回False
            return False

    def _extract_video_frames(self, video_path, num_frames=5, max_dimension=800):
        """从视频中提取帧进行分析"""
        try:
            # 创建临时目录来存储帧
            temp_dir = tempfile.mkdtemp()
            frames = []
            
            # 打开视频文件
            cap = cv2.VideoCapture(video_path)
            if not cap.isOpened():
                if self.progress_callback:
                    self.progress_callback(f"无法打开视频文件：{video_path}")
                return frames, {"width": width, "height": height, "fps": fps, "duration": 0, "frame_count": 0}
            
            # 获取视频信息
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            fps = cap.get(cv2.CAP_PROP_FPS)
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            duration = total_frames / fps if fps > 0 else 0
            
            # 确保至少提取一帧
            if total_frames == 0:
                if self.progress_callback:
                    self.progress_callback(f"视频没有帧：{video_path}")
                return frames, {"width": width, "height": height, "fps": fps, "duration": 0, "frame_count": 0}
            
            # 只提取有限数量的帧
            frame_indices = []
            if total_frames <= num_frames:
                # 如果总帧数小于等于需要的帧数，全部使用
                frame_indices = [i for i in range(total_frames)]
            else:
                # 否则均匀分布选择帧
                for i in range(num_frames):
                    frame_indices.append(int(i * total_frames / num_frames))
            
            if self.progress_callback:
                self.progress_callback(f"从视频中提取 {len(frame_indices)} 帧用于分析: {video_path}")
            
            # 提取并保存帧
            for i, frame_idx in enumerate(frame_indices):
                cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
                ret, frame = cap.read()
                if not ret:
                    continue
                    
                # 调整大小
                if max(height, width) > max_dimension:
                    scale = max_dimension / max(height, width)
                    new_width = int(width * scale)
                    new_height = int(height * scale)
                    frame = cv2.resize(frame, (new_width, new_height))
                
                # OpenCV使用BGR格式，需要转换为RGB
                frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                pil_img = Image.fromarray(frame_rgb)
                
                # 将帧转换为base64字符串
                buffered = BytesIO()
                pil_img.save(buffered, format="JPEG")
                img_str = base64.b64encode(buffered.getvalue()).decode('utf-8')
                frames.append(img_str)
                
                if self.progress_callback:
                    progress_percent = int((i+1) / len(frame_indices) * 100)
                    self.progress_callback(f"已提取 {i+1}/{len(frame_indices)} 帧", progress_percent)
            
            # 释放资源
            cap.release()
            
            # 获取视频元数据
            duration_str = str(datetime.timedelta(seconds=int(duration)))
            metadata = {
                "width": width,
                "height": height,
                "fps": fps,
                "duration": duration,
                "duration_str": duration_str,
                "frame_count": total_frames,
            }
            
            return frames, metadata
            
        except Exception as e:
            if self.progress_callback:
                self.progress_callback(f"提取视频帧失败：{str(e)}")
            print(f"提取视频帧失败：{str(e)}")
            return [], {}

    def _analyze_video(self, video_path):
        """分析视频内容"""
        # 检查是否启用视频分析
        if not self.enable_video_analysis:
            if self.progress_callback:
                self.progress_callback(f"跳过视频分析（未启用）：{video_path}")
            # 返回基本信息而不是完整分析
            file_name = os.path.basename(video_path)
            return {
                "type": "video",
                "content": f"视频文件（未分析）：{file_name}",
                "metadata": {
                    "filename": file_name,
                    "extension": os.path.splitext(file_name)[1].lower(),
                    "analyzed": False
                }
            }
            
        # 检查缓存
        file_md5 = self._calculate_md5(video_path)
        if file_md5 in self.analysis_cache:
            if self.progress_callback:
                self.progress_callback(MessageType.CACHED_ANALYSIS)
            return self.analysis_cache[file_md5]
            
        try:
            if self.progress_callback:
                self.progress_callback(MessageType.ANALYZING)
            
            # 提取视频帧
            frames, metadata = self._extract_video_frames(video_path)
            
            if not frames:
                raise Exception(f"无法从视频中提取帧：{video_path}")
            
            if self.progress_callback:
                self.progress_callback(f"正在分析提取的视频帧...")
            
            # 构建消息内容
            messages = [
                {
                    "role": "system",
                    "content": "根据提供的视频帧分析视频内容，使用不超过60字的一句话描述视频的主要内容和主题"
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": f"这是一个视频的{len(frames)}个关键帧，请根据这些帧分析视频的内容是什么。视频分辨率：{metadata['width']}x{metadata['height']}，时长：{metadata['duration_str']}"
                        }
                    ]
                }
            ]
            
            # 添加每一帧的图像
            for i, frame in enumerate(frames):
                frame_obj = {
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/jpeg;base64,{frame}"
                    }
                }
                messages[1]["content"].append(frame_obj)
            
            # 使用重试机制调用API
            response = self._call_api_with_retry(
                openai.ChatCompletion.create,
                model=self.video_model,
                messages=messages
            )
            
            # 获取分析结果
            result = {
                "type": "video",
                "content": response.choices[0].message.content,
                "metadata": {
                    "filename": os.path.basename(video_path),
                    "extension": os.path.splitext(video_path)[1].lower(),
                    "width": metadata["width"],
                    "height": metadata["height"],
                    "duration": metadata["duration"],
                    "duration_str": metadata["duration_str"],
                    "frame_count": metadata["frame_count"],
                    "fps": metadata["fps"],
                    "analyzed": True,
                    "frames_analyzed": len(frames)
                }
            }
            
            # 更新缓存
            self.analysis_cache[file_md5] = result
            self._save_analysis_cache()
            return result
            
        except Exception as e:
            raise Exception(f"分析视频 {video_path} 失败：{str(e)}")

    def _analyze_file_wrapper(self, file_info):
        """包装文件分析函数，用于线程池"""
        file_path, file_extension = file_info
        try:
            if self.progress_callback:
                with self.progress_lock:
                    self.progress_callback(f"分析文件：{file_path}")
            
            result = self._analyze_file(file_path)
            return file_path, result
        except Exception as e:
            print(f"分析文件出错: {str(e)}")
            return file_path, str(e)
    
    def _analyze_folder_wrapper(self, folder_path):
        """包装文件夹分析函数，用于线程池"""
        try:
            if self.progress_callback:
                with self.progress_lock:
                    self.progress_callback(f"分析文件夹：{folder_path}")
            
            # 获取文件夹的基本信息
            folder_name = os.path.basename(folder_path)
            folder_size = self._get_folder_size(folder_path)
            file_count = self._count_files_in_folder(folder_path)
            
            # 获取文件夹结构和内容分析
            structure = self._analyze_folder_structure(folder_path)
            
            # 使用AI分析文件夹
            folder_analysis = self._analyze_folder_with_ai(folder_path, structure)
            
            # 构建分析结果
            result = {
                "type": "folder",
                "name": folder_name,
                "size": folder_size,
                "file_count": file_count,
                "structure": structure,
                "analysis": folder_analysis,
                "category": folder_analysis.get("category", "其他"),
                "purpose": folder_analysis.get("purpose", "未知"),
                "description": folder_analysis.get("description", f"包含{file_count}个文件的文件夹"),
                "has_irrelevant_files": folder_analysis.get("has_irrelevant_files", False)
            }
            
            return folder_path, result
        except Exception as e:
            print(f"分析文件夹出错: {str(e)}")
            return folder_path, str(e)
    
    def _get_folder_size(self, folder_path):
        """获取文件夹总大小（字节）"""
        total_size = 0
        for dirpath, _, filenames in os.walk(folder_path):
            for f in filenames:
                fp = os.path.join(dirpath, f)
                if os.path.exists(fp):
                    total_size += os.path.getsize(fp)
        return total_size
    
    def _count_files_in_folder(self, folder_path):
        """统计文件夹中的文件数量"""
        count = 0
        for _, _, files in os.walk(folder_path):
            count += len(files)
        return count
    
    def _analyze_folder_structure(self, folder_path, max_depth=4):
        """分析文件夹结构，识别可能格格不入的文件"""
        structure = {
            "files": [],
            "directories": [],
            "extensions": {},
            "has_coherent_structure": False,
            "mismatched_files": []
        }
        
        try:
            # 收集文件夹内的文件和子文件夹
            for root, dirs, files in os.walk(folder_path):
                # 跳过系统文件夹和隐藏文件夹
                dirs[:] = [d for d in dirs if not d.startswith('.') and not d.startswith('$')]
                files = [f for f in files if not f.startswith('.') and not f.startswith('$')]
                
                # 计算当前层级相对于基础文件夹的深度
                depth = root.replace(folder_path, '').count(os.sep)
                if depth > max_depth:
                    continue
                    
                rel_path = os.path.relpath(root, folder_path)
                if rel_path != '.':
                    structure["directories"].append(rel_path)
                
                for file in files:
                    file_path = os.path.relpath(os.path.join(root, file), folder_path)
                    structure["files"].append(file_path)
                    
                    # 统计文件扩展名
                    ext = os.path.splitext(file)[1].lower()
                    if ext:
                        structure["extensions"][ext] = structure["extensions"].get(ext, 0) + 1
        except Exception as e:
            self.logger.log_error(f"分析文件夹结构失败：{str(e)}")
            return structure
        
        # 分析文件类型一致性
        if len(structure["extensions"]) <= 2:  # 如果只有1-2种文件类型
            # 检查是否存在明显的主要文件类型
            total_files = len(structure["files"])
            threshold = 0.8  # 80%的文件属于同一类型，认为结构一致
            
            for ext, count in structure["extensions"].items():
                if count / total_files >= threshold:
                    structure["has_coherent_structure"] = True
                    
                    # 找出不符合主要扩展名的文件
                    for file in structure["files"]:
                        file_ext = os.path.splitext(file)[1].lower()
                        if file_ext != ext:
                            structure["mismatched_files"].append(file)
                    
                    break
        
        # 如果没有找到主要文件类型，尝试用AI分析
        if not structure["has_coherent_structure"] and len(structure["files"]) > 5:
            # 在实际使用中，这里会调用AI分析文件夹结构
            # 简化版：如果有统一的文件命名模式，认为结构一致
            file_names = [os.path.splitext(os.path.basename(f))[0] for f in structure["files"]]
            if self._has_naming_pattern(file_names):
                structure["has_coherent_structure"] = True
        
        return structure
    
    def _has_naming_pattern(self, file_names):
        """检查文件名是否有统一的命名模式"""
        if len(file_names) < 3:
            return False
            
        # 简单检查：是否有数字序列
        numbered_files = 0
        for name in file_names:
            if re.search(r'\d+', name):
                numbered_files += 1
        
        return numbered_files / len(file_names) >= 0.7  # 70%以上的文件有数字
    
    def _analyze_folder_with_ai(self, folder_path, structure):
        """使用AI分析文件夹内容和结构"""
        try:
            # 如果文件夹中的文件很少，不需要AI分析
            if len(structure["files"]) <= 3:
                return {
                    "purpose": "未知",
                    "category": "其他",
                    "description": f"包含{len(structure['files'])}个文件的文件夹"
                }
            
            # 获取文件夹名和一些示例文件
            folder_name = os.path.basename(folder_path)
            sample_files = structure["files"][:5]  # 最多取5个文件作为示例
            
            # 构建提示词
            extensions_summary = ", ".join([f"{ext} ({count}个)" for ext, count in structure["extensions"].items()])
            prompt = f"""分析这个文件夹的结构和内容，并确定其用途和类别：
文件夹名称: {folder_name}
文件数量: {len(structure['files'])}
子文件夹数量: {len(structure['directories'])}
文件类型统计: {extensions_summary}
示例文件: {', '.join(sample_files)}

请提供：
1. 这个文件夹的主要用途(purpose)
2. 文件夹的最合适分类(category)，可参考选项：文档、图片、视频、音乐、代码、数据库、归档/备份、应用程序、系统、其他
3. 文件夹内容的简短描述(description，50字以内)
4. 是否包含无关文件(has_irrelevant_files)，返回true或false"""

            # 调用AI分析
            response = self._call_api_with_retry(
                openai.ChatCompletion.create,
                model=self.file_model,
                messages=[{
                    "role": "user",
                    "content": prompt
                }]
            )
            
            # 解析AI返回结果
            analysis = self._parse_ai_folder_analysis(response)
            
            return analysis
            
        except Exception as e:
            print(f"AI分析文件夹失败: {str(e)}")
            return {
                "purpose": "未知",
                "category": "其他",
                "description": "AI分析失败"
            }
    
    def _parse_ai_folder_analysis(self, ai_response):
        """解析AI返回的文件夹分析结果"""
        try:
            # 从API响应中提取文本内容
            if hasattr(ai_response, 'choices') and len(ai_response.choices) > 0:
                content = ai_response.choices[0].message.content
            else:
                content = str(ai_response)  # 如果不是标准响应对象，尝试转为字符串
            
            # 尝试解析结构化数据
            if '{' in content and '}' in content:
                # 尝试提取JSON部分
                import re
                json_match = re.search(r'{.*}', content, re.DOTALL)
                if json_match:
                    json_str = json_match.group(0)
                    result = json.loads(json_str)
                    return result
            
            # 如果不是结构化数据，尝试从文本中提取关键信息
            import re  # 确保re模块在这里也能被引用
            analysis = {
                "purpose": "未知",
                "category": "其他",
                "description": "文件夹",
                "has_irrelevant_files": False
            }
            
            # 尝试提取purpose
            purpose_match = re.search(r'用途[：:]\s*(.+?)[,\n]', content)
            if purpose_match:
                analysis["purpose"] = purpose_match.group(1).strip()
                
            # 尝试提取category
            categories = ["文档", "图片", "视频", "音乐", "代码", "数据库", "归档", "备份", "应用程序", "系统", "其他"]
            for cat in categories:
                if cat in content:
                    analysis["category"] = cat
                    break
            
            # 尝试提取description
            desc_match = re.search(r'描述[：:]\s*(.+?)[,\n]', content)
            if desc_match:
                analysis["description"] = desc_match.group(1).strip()
            
            # 检查是否有无关文件
            if "无关文件" in content and ("是" in content or "true" in content.lower()):
                analysis["has_irrelevant_files"] = True
                
            return analysis
            
        except Exception as e:
            self.logger.log_error(f"解析AI返回的文件夹分析结果失败: {str(e)}")
            return {
                "purpose": "未知",
                "category": "其他",
                "description": "解析AI分析结果失败"
            }

    def cancel_operation(self):
        """设置取消标志，中断当前执行的操作"""
        self.cancel_flag = True
        if self.progress_callback:
            self.progress_callback("正在取消操作，请稍等...")
        self.logger.log_info("用户取消操作")
        
    def reset_cancel_flag(self):
        """重置取消标志"""
        self.cancel_flag = False

    def organize_directory(self, directory_path):
        self.reset_cancel_flag()  # 重置取消标志
        analysis_results = {}
        file_list = []
        
        # 获取配置的线程数
        try:
            thread_count = self.config.getint('Settings', 'thread_count', fallback=8)
            thread_count = max(1, min(32, thread_count))  # 确保线程数在1-32之间
        except:
            thread_count = 8
        
        # 根据子文件夹处理模式收集需要分析的文件
        if self.subfolder_mode == 'extract_all':
            # 全部解体模式：收集所有文件
            for root, dirs, files in os.walk(directory_path):
                if self.cancel_flag:  # 检查取消标志
                    if self.progress_callback:
                        self.progress_callback("操作已取消")
                    return None
                for file_name in files:
                    file_path = os.path.join(root, file_name)
                    file_extension = os.path.splitext(file_name)[1].lower()
                    file_list.append((file_path, file_extension))
        elif self.subfolder_mode == 'whole':
            # 整体处理模式：只收集顶级文件和文件夹
            for item in os.listdir(directory_path):
                if self.cancel_flag:  # 检查取消标志
                    if self.progress_callback:
                        self.progress_callback("操作已取消")
                    return None
                item_path = os.path.join(directory_path, item)
                if os.path.isfile(item_path):
                    file_extension = os.path.splitext(item)[1].lower()
                    file_list.append((item_path, file_extension))
                elif os.path.isdir(item_path):
                    # 将文件夹作为整体添加到分析列表
                    file_list.append((item_path, 'folder'))
        else:  # subfolder_mode == 'extract_partial'
            # 部分提取模式：需要先分析每个文件夹
            for item in os.listdir(directory_path):
                if self.cancel_flag:  # 检查取消标志
                    if self.progress_callback:
                        self.progress_callback("操作已取消")
                    return None
                item_path = os.path.join(directory_path, item)
                if os.path.isfile(item_path):
                    # 单独文件直接添加
                    file_extension = os.path.splitext(item)[1].lower()
                    file_list.append((item_path, file_extension))
                elif os.path.isdir(item_path):
                    # 分析文件夹
                    folder_analysis = self._analyze_folder_structure(item_path)
                    if folder_analysis.get("mismatched_files"):
                        # 如果有格格不入的文件，将它们单独添加
                        for mismatched_file in folder_analysis["mismatched_files"]:
                            file_path = os.path.join(item_path, mismatched_file)
                            file_extension = os.path.splitext(mismatched_file)[1].lower()
                            file_list.append((file_path, file_extension))
                    
                    if folder_analysis.get("has_coherent_structure") and not folder_analysis.get("mismatched_files"):
                        # 如果文件夹结构一致且没有格格不入的文件，将文件夹作为整体添加
                        file_list.append((item_path, 'folder'))
                    elif not folder_analysis.get("mismatched_files"):
                        # 如果没有格格不入的文件但结构不一致，也将文件夹作为整体添加
                        file_list.append((item_path, 'folder'))
        
        # 使用线程池并行处理文件
        with ThreadPoolExecutor(max_workers=thread_count) as executor:
            # 提交所有任务
            future_to_file = {}
            for file_info in file_list:
                if self.cancel_flag:  # 检查取消标志
                    # 取消已经提交的任务
                    for future in future_to_file:
                        future.cancel()
                    if self.progress_callback:
                        self.progress_callback("操作已取消")
                    return None
                
                file_path, file_extension = file_info
                if file_extension == 'folder':
                    # 处理文件夹
                    future = executor.submit(self._analyze_folder_wrapper, file_path)
                else:
                    # 处理文件
                    future = executor.submit(self._analyze_file_wrapper, file_info)
                future_to_file[future] = file_path
            
            # 收集结果
            for future in future_to_file:
                if self.cancel_flag:  # 检查取消标志
                    # 尝试取消剩余的任务
                    for f in future_to_file:
                        if not f.done():
                            f.cancel()
                    if self.progress_callback:
                        self.progress_callback("操作已取消")
                    return None
                
                file_path, result = future.result()
                if isinstance(result, str):  # 如果是错误信息
                    if self.progress_callback:
                        with self.progress_lock:
                            self.progress_callback(f"分析失败：{file_path}，错误：{result}")
                    continue
                analysis_results[file_path] = result

        if self.cancel_flag:  # 再次检查取消标志
            if self.progress_callback:
                self.progress_callback("操作已取消")
            return None

        # 获取最终整理决策
        decision = self._get_final_decision(analysis_results)
        
        if self.cancel_flag or not decision:  # 检查取消标志或决策是否成功
            if self.progress_callback:
                self.progress_callback("操作已取消或生成决策失败")
            return None

        # 处理目标路径和文件名
        if decision and 'files' in decision:
            # 创建目标路径映射，检测重复
            path_map = {}  # 记录每个目标路径及其对应的文件列表
            folder_files = {}  # 记录每个目标文件夹中的文件名
            
            # 第一遍：收集所有目标路径和文件名，并验证路径安全性
            invalid_paths = []  # 记录不合法的路径
            for file_info in decision['files']:
                src = file_info['original_path']
                dst = file_info['new_path']
                
                # 保持原始文件名
                is_folder = os.path.isdir(src)
                original_name = os.path.basename(src)
                target_folder = os.path.dirname(dst)
                
                # 构建新的目标路径（使用原始文件名/文件夹名）
                new_dst = os.path.join(target_folder, original_name)
                file_info['new_path'] = new_dst
                
                # 计算完整目标路径
                if self.output_dir:
                    target_path = os.path.join(self.output_dir, new_dst.lstrip(os.sep))
                else:
                    target_path = os.path.join(os.path.dirname(src), new_dst.lstrip(os.sep))
                
                # 验证目标路径是否合法
                if not self._validate_output_path(target_path):
                    invalid_paths.append((src, target_path))
                    continue
                
                # 记录目标路径和源文件
                if target_path not in path_map:
                    path_map[target_path] = []
                path_map[target_path].append(src)
                
                # 记录目标文件夹中的文件/文件夹
                target_dir = os.path.dirname(target_path)
                if target_dir not in folder_files:
                    folder_files[target_dir] = set()
                folder_files[target_dir].add(original_name)
                
                # 标记是文件夹还是文件
                file_info['is_folder'] = is_folder
            
            # 如果有不合法的路径，返回错误
            if invalid_paths:
                error_msg = "发现不合法的目标路径：\n"
                for src, dst in invalid_paths:
                    error_msg += f"- {src} -> {dst}\n"
                if self.progress_callback:
                    self.progress_callback(error_msg)
                print(error_msg)
                return None
            
            # 第二遍：处理文件名冲突
            for file_info in decision['files']:
                src = file_info['original_path']
                dst = file_info['new_path']
                
                # 计算完整目标路径
                if self.output_dir:
                    target_path = os.path.join(self.output_dir, dst.lstrip(os.sep))
                else:
                    target_path = os.path.join(os.path.dirname(src), dst.lstrip(os.sep))
                
                # 检查是否存在路径冲突
                if len(path_map[target_path]) > 1:
                    # 获取文件名和扩展名
                    file_dir = os.path.dirname(dst)
                    original_filename = os.path.basename(src)
                    name, ext = os.path.splitext(original_filename)
                    
                    # 确定当前文件在冲突列表中的索引
                    conflict_index = path_map[target_path].index(src)
                    
                    # 只有索引大于0的文件需要重命名（第一个文件保持原名）
                    if conflict_index > 0:
                        # 创建新的路径，添加序号
                        new_name = f"{name}_{conflict_index+1}{ext}"
                        new_dst = os.path.join(file_dir, new_name)
                        
                        # 更新决策中的路径
                        file_info['new_path'] = new_dst
                        # 记录原始目标路径，用于还原
                        file_info['original_target'] = dst
                        
                        # 更新目标路径映射
                        if self.output_dir:
                            new_target_path = os.path.join(self.output_dir, new_dst.lstrip(os.sep))
                        else:
                            new_target_path = os.path.join(os.path.dirname(src), new_dst.lstrip(os.sep))
                        
                        # 再次验证新路径的安全性
                        if not self._validate_output_path(new_target_path):
                            error_msg = f"发现不合法的目标路径：\n- {src} -> {new_target_path}"
                            if self.progress_callback:
                                self.progress_callback(error_msg)
                            print(error_msg)
                            return None
                        
                        if new_target_path not in path_map:
                            path_map[new_target_path] = []
                        path_map[new_target_path].append(src)
        
        # 返回决策供GUI确认
        return decision

    def restore_files(self):
        # 实现还原功能
        if not os.path.exists(self.history_file):
            print("没有发现恢复的历史")
            return None
        
        restore_plan = []
        for original_path, info in self.history.items():
            if info["new_path"] and os.path.exists(info["new_path"]):
                # 直接还原到原始路径
                restore_path = original_path
                
                restore_plan.append({
                    "current_path": info["new_path"],
                    "restore_path": restore_path,
                    "original_path": original_path  # 保留原始路径信息用于GUI显示
                })
        
        # 返回还原计划供界面确认
        return {"files": restore_plan} if restore_plan else None
